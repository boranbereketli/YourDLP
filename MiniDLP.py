#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import time
import os
import shutil
import pyperclip
import threading
import re
import socket
import platform
import ctypes
import string
from dataclasses import dataclass
from typing import Optional
from types import SimpleNamespace
import requests
import json


from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# Document libs
from docx import Document
import PyPDF2
import pandas as pd
from pptx import Presentation


# ============================================================
# CONFIG
# ============================================================

LOG_CSV = "dlp_incidents.csv"
QUARANTINE_DIR = "KARANTINA_ALANI"
os.makedirs(QUARANTINE_DIR, exist_ok=True)

SIM_USB_DIR = "SIM_USB_SURUCU"
os.makedirs(SIM_USB_DIR, exist_ok=True)

MAX_FILE_SIZE = 15 * 1024 * 1024  # 15 MB
ALLOWED_EXT = {".txt", ".csv", ".docx", ".pdf", ".xlsx", ".xls", ".pptx"}

CONFIG_FILE = "dlp_config.json"

# Varsayılan SCAN_RULES
SCAN_RULES = {
    "TCKN": True,
    "TEL_NO": True,
    "IBAN_TR": True,
    "KREDI_KARTI": True,
    "E_POSTA": True
}

# Config dosyası varsa üzerinden güncelle
if os.path.exists(CONFIG_FILE):
    try:
        with open(CONFIG_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            SCAN_RULES.update(data.get("SCAN_RULES", {}))
        print("[Config] dlp_config.json yüklendi.")
    except Exception:
        print("[Config] dlp_config.json okunamadı, varsayılan ayarlar kullanılıyor.")
else:
    # Yoksa otomatik oluştur
    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
        json.dump({"SCAN_RULES": SCAN_RULES}, f, indent=4)
    print("[Config] Varsayılan dlp_config.json oluşturuldu.")


# ============================================================
# USB POLICY MODE (STRICT: tüm dosyaları karantinaya alır, SMART: sadece hassas veri içerenleri karantinaya alır)
# ============================================================

USB_POLICY = "SMART"  # Varsayılan ayar

# ============================================================
# DYNAMIC BANNED WORDS (regex değil, düz kelime arama)
# ============================================================

BANNED_WORDS = set()  # runtime'da eklenecek

# ============================================================
# HELPERS: TCKN / PHONE / IBAN VALIDATORS
# ============================================================

# ==== Policy Server ile entegre değişkenler ====

USB_SCAN_ON_MODIFY = True     # USB değişiklikte taransın mı?

CLIPBOARD_ENABLED = True      # True / False
USB_ENABLED       = True      # True / False
NETWORK_ENABLED   = True      # True / False

# Policy Server URL
POLICY_URL = "http://127.0.0.1:5000/get_policy/PC1"
# ============================================================

def is_valid_tckn(tckn: str) -> bool:
    """
    Türkiye Cumhuriyeti Kimlik Numarası doğrulaması (checksum kuralları).
    """
    if not isinstance(tckn, str):
        return False
    t = re.sub(r"\D", "", tckn)
    if len(t) != 11:
        return False
    if t[0] == "0":
        return False
    try:
        digits = list(map(int, t))
    except ValueError:
        return False

    rule_10 = ((digits[0] + digits[2] + digits[4] + digits[6] + digits[8]) * 7 - 
               (digits[1] + digits[3] + digits[5] + digits[7])) % 10
    if rule_10 != digits[9]:
        return False

    rule_11 = sum(digits[:10]) % 10
    if rule_11 != digits[10]:
        return False

    return True


def is_valid_phone(phone: str) -> bool:
    """
    Türkiye mobil telefon doğrulaması.
    Kabul edilen formatlar:
      - 05XXXXXXXXX (11 hane, başında 0)
      - +905XXXXXXXXX
      - 5XXXXXXXXX (10 hane, başında 0 eksik)
      - 90XXXXXXXXXX (12 hane, başında 90)
    Dönen: True/False
    """
    if not isinstance(phone, str):
        return False
    digits = re.sub(r"\D", "", phone)
    # +90 prefix case -> '90' + 10 digits => length 12 and startswith '90'
    if digits.startswith("90") and len(digits) == 12:
        national = digits[2:]
    # leading 0 case -> '0' + 10 digits => len 11 and startswith '0'
    elif digits.startswith("0") and len(digits) == 11:
        national = digits[1:]
    # maybe given without leading 0 -> 10 digits starting with 5
    elif len(digits) == 10:
        national = digits
    else:
        return False

    # Now national should be 10 digits, starting with '5' for mobile numbers in TR
    if len(national) != 10:
        return False
    if not national.isdigit():
        return False
    if not national.startswith("5"):
        return False
    return True


def iban_to_numeric(iban: str) -> Optional[str]:
    """
    IBAN string -> numeric string for mod-97
    Example: 'TR33...' -> move first 4 chars, expand letters A=10..Z=35
    """
    try:
        s = re.sub(r"\s+", "", iban).upper()
        if len(s) < 4:
            return None
        rearr = s[4:] + s[:4]
        numeric = []
        for ch in rearr:
            if ch.isdigit():
                numeric.append(ch)
            elif ch.isalpha():
                numeric.append(str(ord(ch) - 55))  # A->10
            else:
                return None
        return "".join(numeric)
    except Exception:
        return None


def is_valid_iban(iban: str) -> bool:
    """
    IBAN mod-97 validation. Returns True if valid.
    """
    if not isinstance(iban, str):
        return False
    s = re.sub(r"\s+", "", iban).upper()
    if not re.match(r"^[A-Z]{2}\d{2}[A-Z0-9]+$", s):
        return False
    numeric = iban_to_numeric(s)
    if not numeric:
        return False
    # iterative mod-97 to avoid huge ints
    remainder = 0
    chunk_size = 9
    for i in range(0, len(numeric), chunk_size):
        chunk = str(remainder) + numeric[i:i + chunk_size]
        remainder = int(chunk) % 97
    return remainder == 1


# ============================================================
# DLP RULES (güncellendi: TEL_NO, TCKN, IBAN_TR)
# ============================================================

TEL_REGEX = r'(?:(?:\+90|0)?5\d{9})'  # matches +905xxxxxxxxx or 05xxxxxxxxx or 5xxxxxxxxx
DLP_RULES = {
    "TCKN": {
        "pattern": r'\b[1-9]\d{10}\b',  # 11 digits not starting with 0 (candidate)
        "description": "11 Haneli TC Kimlik Numarası"
    },
    "TEL_NO": {
        "pattern": TEL_REGEX,
        "description": "Türkiye Telefon Numarası"
    },
    "KREDI_KARTI": {
        "pattern": r'\b\d{4}[- ]?\d{4}[- ]?\d{4}[- ]?\d{4}\b',
        "description": "16 Haneli Kredi Kartı Numarası Formatı"
    },
    "E_POSTA": {
        "pattern": r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}',
        "description": "E-posta Adresi Formatı"
    },
    "IBAN_TR": {
        # Non-capturing groups to avoid tuple matches
        "pattern": r'\bTR\d{2}[A-Z0-9]{4}\s?(?:\d{4}\s?){4}\d{2}\b',
        "description": "Türk IBAN Numarası Formatı"
    }
}

# ==========================
# MONITOR STARTER FUNCTIONS
# ==========================

def start_clipboard_monitor():
    t = threading.Thread(target=clipboard_monitor, daemon=True)
    t.start()
    print("[DLP] Clipboard monitor başlatıldı.")

def start_usb_monitor():
    t = threading.Thread(target=usb_monitor, daemon=True)
    t.start()
    print("[DLP] USB monitor başlatıldı.")

def start_network_monitor():
    print("[DLP] Network DLP henüz entegre edilmedi (stub).")
    # ileride: sender + gateway + receiver otomatik başlatılabilir


def scan_content(content: str):
    """
    İçeriği tarar ve tespit edilen hassas verileri döndürür.
    DÖNÜŞ FORMAT:
    [
        {"data_type": "TCKN", "masked_match": "TC: ******1234"},
        {"data_type": "BANNED_WORD", "masked_match": "password"},
        ...
    ]
    """

    incidents = []  # SONUÇ LİSTESİ

    if not content:
        return incidents

    text = str(content)

    # =========================================================
    # 1) TELEFON EN ÖNCE – karışmasın diye
    # =========================================================
    tel_regex = re.compile(DLP_RULES["TEL_NO"]["pattern"])
    for match in tel_regex.findall(text):
        if is_valid_phone(match):
            flat = re.sub(r"\D", "", match)
            incidents.append({
                "data_type": "TEL_NO",
                "masked_match": f"TEL: ******{flat[-2:]}"
            })
    # Telefonları TCKN taramasından çıkarmak için çıkarıyoruz
    cleaned_text = tel_regex.sub(" ", text)

    # =========================================================
    # 2) TCKN TESPITI
    # =========================================================
    tckn_regex = re.compile(DLP_RULES["TCKN"]["pattern"])
    for match in tckn_regex.findall(cleaned_text):
        only_digits = re.sub(r"\D", "", match)
        if is_valid_tckn(only_digits):
            incidents.append({
                "data_type": "TCKN",
                "masked_match": f"TC: ******{only_digits[-4:]}"
            })

    # =========================================================
    # 3) IBAN, CC, EMAIL — NORMAL REGEX TARAMASI
    # =========================================================
    for dtype, rule in DLP_RULES.items():
        if not SCAN_RULES.get(dtype, True):   # pasifse bu rule atlanır
            continue
        if dtype in ("TEL_NO", "TCKN"):
            continue
  
        regex = re.compile(rule["pattern"])
        for match in regex.findall(text):
            m = match if isinstance(match, str) else "".join(match)

            if dtype == "IBAN_TR" and is_valid_iban(m):
                incidents.append({
                    "data_type": "IBAN_TR",
                    "masked_match": f"IBAN: ****{m[-4:]}"
                })

            elif dtype == "KREDI_KARTI":
                digits = re.sub(r"\D", "", m)
                incidents.append({
                    "data_type": "KREDI_KARTI",
                    "masked_match": f"CC: XXXX...{digits[-4:]}"
                })

            elif dtype == "E_POSTA":
                incidents.append({
                    "data_type": "E_POSTA",
                    "masked_match": f"EMAIL: <{m.split('@')[0][0]}***@...>"
                })

    # =========================================================
    # 4) YASAKLI KEYWORD TARAMASI (EXTRA)
    # =========================================================
    for word in BANNED_WORDS:
        if word.lower() in text.lower():
            incidents.append({
                "data_type": "BANNED_WORD",
                "masked_match": word
            })

    return incidents


def add_banned_word(word: str):
    BANNED_WORDS.add(word.strip().lower())
    print(f"[DLP] Yasaklı kelime eklendi: {word}")


def remove_banned_word(word: str):
    """Yasaklı kelime listeden çıkarılır. Yoksa uyarı verir."""
    cleaned = word.strip().lower()
    if cleaned in BANNED_WORDS:
        BANNED_WORDS.remove(cleaned)
        print(f"[DLP] Yasaklı kelime kaldırıldı: {word}")
    else:
        print(f"[DLP] Kelime bulunamadı: {word}")


def list_banned_words():
    """Şu anda aktif yasaklı kelimeleri gösterir."""
    if not BANNED_WORDS:
        print("[DLP] Henüz yasaklı kelime eklenmemiş.")
    else:
        print("\n--- AKTİF YASAKLI KELİMELER ---")
        for w in BANNED_WORDS:
            print(f" - {w}")
        print("-------------------------------\n")



# ============================================================
# LOGGING
# ============================================================

def log_incident(event_type, data_type, action, details):
    log_line = f"{time.strftime('%Y-%m-%d %H:%M:%S')},{event_type},{data_type},{action},{details}\n"
    try:
        if not os.path.exists(LOG_CSV):
            with open(LOG_CSV, "w", encoding="utf-8") as f:
                f.write("Tarih,Olay_Tipi,Veri_Tipi,Aksiyon,Detay\n")
        with open(LOG_CSV, "a", encoding="utf-8") as f:
            f.write(log_line)
        print(f"\n[LOG] {data_type} | {action} | {details}")
    except PermissionError:
        print("[LOG ERROR] PermissionError while writing log.")
    except Exception as e:
        print(f"[LOG ERROR] {e}")


# ============================================================
# FILE READING (TXT, DOCX, PDF, XLSX, PPTX)
# ============================================================

def read_file_content(path: str) -> str:
    try:
        ext = os.path.splitext(path)[1].lower()
    except Exception:
        return ""

    try:
        # size guard
        try:
            if os.path.getsize(path) > MAX_FILE_SIZE:
                print(f"[READ] Skipping large file: {path}")
                return ""
        except Exception:
            pass

        if ext in (".txt", ".csv"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        if ext == ".docx":
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        if ext == ".pdf":
            text = ""
            with open(path, "rb") as f:
                try:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() or ""
                except Exception:
                    return ""
            return text

        if ext in (".xlsx", ".xls"):
            try:
                dfs = pd.read_excel(path, sheet_name=None)
                return "\n".join(df.to_string(index=False) for df in dfs.values())
            except Exception:
                return ""

        if ext == ".pptx":
            try:
                prs = Presentation(path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            parts.append(shape.text)
                return "\n".join(parts)
            except Exception:
                return ""

        return ""
    except Exception as e:
        print(f"[READ ERROR] {e}")
        return ""


# ============================================================
# USB HANDLER & MONITOR
# ============================================================

class USBFileHandler(FileSystemEventHandler):
    def __init__(self, quarantine_dir=QUARANTINE_DIR):
        super().__init__()
        self.quarantine_dir = quarantine_dir

    def _quarantine(self, src_path, hint_name=None):
        name = hint_name or os.path.basename(src_path)
        dest = os.path.join(self.quarantine_dir, f"{int(time.time())}_{name}")
        try:
            shutil.move(src_path, dest)
            return dest
        except Exception:
            try:
                shutil.copy2(src_path, dest)
                try:
                    os.remove(src_path)
                except Exception:
                    pass
                return dest
            except Exception:
                return None

    def _process_file(self, file_path):
        try:
            if not os.path.isfile(file_path):
                return
            name = os.path.basename(file_path)
            if name.startswith(".") or name.startswith("~$"):
                return
            try:
                size = os.path.getsize(file_path)
            except Exception:
                size = 0
            if size == 0:
                return

            if size > MAX_FILE_SIZE:
                print(f"[USB] Skipping large file: {file_path}")
                return

            ext = os.path.splitext(name)[1].lower()
            content = ""
            if ext in ALLOWED_EXT:
                content = read_file_content(file_path)

            incidents = scan_content(content) if content else []
            if incidents:
                types = ", ".join({i['data_type'] for i in incidents})
                masked = incidents[0]['masked_match']
                details = f"{name} -> {masked}"
                log_incident("USB Transferi", types, "ENGEL - Karantina", details)
                q = self._quarantine(file_path, hint_name=name)
                if q:
                    print(f"!!! Dosya karantinaya taşındı: {q} !!!")
                else:
                    print(f"[USB] Karantina başarısız: {file_path}")
            else:
                if USB_POLICY == "SMART":
                    log_incident("USB Yazma", "TEMIZ", "İZİN VERİLDİ", name)
                    print(f"[SMART MODE] Temiz dosya USB'ye yazıldı: {name}")
                    return   

                else:
                    q = self._quarantine(file_path, hint_name=name)
                    if q:
                        log_incident("USB Yazma", "Bilinmiyor", "ENGEL - Karantinaya alındı", name)
                        print(f"[STRICT MODE] {name} karantinaya alındı (temiz olsa bile).")
                    return
                if q:
                    log_incident("USB Yazma", "Bilinmiyor", "ENGEL - Karantinaya alındı", name)
                    print(f"USB'ye yazma politikası: {name} karantinaya alındı (temiz olsa bile).")
                else:
                    print(f"[USB] Karantinaya alma başarısız (temiz dosya): {file_path}")

        except Exception as e:
            log_incident("USB Transferi", "Hata", "İşlenmedi", f"Dosya işleme hatası: {e}")

    def on_created(self, event):
        if event.is_directory:
            return
        time.sleep(0.1)
        self._process_file(event.src_path)

    def on_modified(self, event):
        if event.is_directory:
            return
        time.sleep(0.05)
        self._process_file(event.src_path)


def get_usb_mount_points():
    mounts = []
    system = platform.system()
    try:
        if system == "Windows":
            DRIVE_REMOVABLE = 2
            bitmask = ctypes.windll.kernel32.GetLogicalDrives()
            for i, letter in enumerate(string.ascii_uppercase):
                if bitmask & (1 << i):
                    drive = f"{letter}:\\"
                    try:
                        drive_type = ctypes.windll.kernel32.GetDriveTypeW(ctypes.c_wchar_p(drive))
                        if drive_type == DRIVE_REMOVABLE:
                            mounts.append(drive)
                    except Exception:
                        continue
        else:
            base_dirs = ["/media", "/run/media"]
            for base in base_dirs:
                if os.path.exists(base):
                    for entry in os.listdir(base):
                        candidate = os.path.join(base, entry)
                        if os.path.ismount(candidate) or os.path.isdir(candidate):
                            mounts.append(candidate)
                            for sub in os.listdir(candidate):
                                subpath = os.path.join(candidate, sub)
                                if os.path.ismount(subpath) or os.path.isdir(subpath):
                                    mounts.append(subpath)
    except Exception:
        pass

    if os.path.exists(SIM_USB_DIR) and SIM_USB_DIR not in mounts:
        mounts.append(SIM_USB_DIR)

    unique = []
    for m in mounts:
        if m and m not in unique:
            unique.append(m)
    return unique


def scan_existing_files_in_mount(mount_path):
    handler = USBFileHandler()
    for root, _, files in os.walk(mount_path):
        for f in files:
            fp = os.path.join(root, f)
            try:
                fake_event = SimpleNamespace(is_directory=False, src_path=fp)
                handler.on_created(fake_event)
            except Exception:
                continue


def start_observer_for_mount(mount_path):
    try:
        scan_existing_files_in_mount(mount_path)
        event_handler = USBFileHandler()
        observer = Observer()
        observer.schedule(event_handler, mount_path, recursive=True)
        observer.daemon = True
        observer.start()
        print(f"[USB OBSERVER] Başlatıldı: {mount_path}")
        return observer
    except Exception as e:
        print(f"[USB OBSERVER] Başlatılamadı ({mount_path}): {e}")
        return None


def usb_monitor():
    print("[USB] Gerçek USB mount izleyici başlatıldı.")
    known = set()
    observers = {}
    if os.path.exists(SIM_USB_DIR):
        obs = start_observer_for_mount(SIM_USB_DIR)
        if obs:
            observers[SIM_USB_DIR] = obs
            known.add(SIM_USB_DIR)

    try:
        while True:
            mounts = set(get_usb_mount_points())
            added = mounts - known
            removed = known - mounts

            for m in added:
                print(f"[USB] Yeni mount bulundu: {m}")
                obs = start_observer_for_mount(m)
                if obs:
                    observers[m] = obs
                    known.add(m)
                else:
                    try:
                        scan_existing_files_in_mount(m)
                        known.add(m)
                    except Exception:
                        pass

            for m in removed:
                print(f"[USB] Mount çıkarıldı: {m}")
                if m in observers:
                    try:
                        observers[m].stop()
                        observers[m].join(timeout=1)
                    except Exception:
                        pass
                    del observers[m]
                known.discard(m)

            time.sleep(2)
    except KeyboardInterrupt:
        for o in observers.values():
            try:
                o.stop()
                o.join()
            except Exception:
                pass
        print("[USB] İzleyici durduruldu.")


# ============================================================
# CLIPBOARD MONITOR
# ============================================================

def clipboard_monitor():
    print("[CLIPBOARD] Pano İzleyici Aktif.")
    last_clipboard_content = None
    while True:
        try:
            try:
                current_clipboard_content = pyperclip.paste()
            except Exception:
                time.sleep(1)
                continue

            if current_clipboard_content is None:
                current_clipboard_content = ""

            if current_clipboard_content != last_clipboard_content and current_clipboard_content:
                incidents = scan_content(str(current_clipboard_content))

                if incidents:
                    detected_types = {i['data_type'] for i in incidents}

                    if len(detected_types) >= 2:
                        risk_level = "YÜKSEK RİSK - Çoklu Veri Sızıntısı"
                    else:
                        risk_level = "ENGEL - Pano Temizlendi"

                    data_type = ", ".join(detected_types)
                    masked_match = incidents[0]['masked_match']

                    log_incident("Pano Kopyalama", data_type, risk_level, masked_match)

                    try:
                        pyperclip.copy("Bu içerik hassas veri içerdiği için DLP tarafından engellenmiştir.")
                    except Exception:
                        pass

                last_clipboard_content = current_clipboard_content

            time.sleep(0.5)
        except Exception:
            time.sleep(1)


# ============================================================
# NETWORK AGENT (gateway / sender / receiver) - unchanged
# ============================================================

@dataclass
class Message:
    src: str
    dst: str
    channel: str
    payload: str


class DLPAgentGateway:
    def __init__(self, name="DLP_GATEWAY"):
        self.name = name

    def handle(self, msg: Message) -> Optional[Message]:
        print(f"\n[{self.name}] Mesaj alındı: {msg.src} -> {msg.dst}")
        print(f"Kanal : {msg.channel}")
        print(f"İçerik: {msg.payload}")

        incidents = scan_content(msg.payload)

        if incidents:
            detected_types = sorted({i["data_type"] for i in incidents})
            data_type_str = "/".join(detected_types)
            masked_samples = ", ".join(sorted({i["masked_match"] for i in incidents}))

            log_incident(
                event_type=f"{msg.channel} Mesajı",
                data_type=data_type_str,
                action="ENGEL - Mesaj gönderilmedi",
                details=f"{msg.src}->{msg.dst} | {masked_samples}",
            )

            print(f"[{self.name}] UYARI: Mesaj BLOKLANDI!")
            print(f"  Veri tipleri : {data_type_str}")
            print(f"  Örnek       : {masked_samples}")
            return None
        else:
            log_incident(
                event_type=f"{msg.channel} Mesajı",
                data_type="YOK",
                action="İZİN VERİLDİ - Mesaj iletildi",
                details=f"{msg.src}->{msg.dst} | {msg.payload[:50]}",
            )

            print(f"[{self.name}] Mesaj temiz, {msg.dst}'ye iletiliyor.")
            return msg


# Network config (fill in if using)
GATEWAY_LISTEN_HOST = "127.0.0.1"
GATEWAY_LISTEN_PORT = 9101
RECEIVER_HOST = "127.0.0.1"
RECEIVER_PORT = 9102
SENDER_GATEWAY_HOST = "127.0.0.1"
SENDER_GATEWAY_PORT = 9101
RECEIVER_LISTEN_HOST = "127.0.0.1"
RECEIVER_LISTEN_PORT = 9102


def run_gateway():
    dlp = DLPAgentGateway()

    print(f"[{dlp.name}] Receiver'a bağlanılıyor: {RECEIVER_HOST}:{RECEIVER_PORT}")
    receiver_sock = socket.create_connection((RECEIVER_HOST, RECEIVER_PORT))
    print(f"[{dlp.name}] Receiver bağlantısı OK.")

    server_sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    server_sock.bind((GATEWAY_LISTEN_HOST, GATEWAY_LISTEN_PORT))
    server_sock.listen(1)
    print(f"[{dlp.name}] Sender için dinleniyor: {GATEWAY_LISTEN_HOST}:{GATEWAY_LISTEN_PORT}")

    sender_sock, sender_addr = server_sock.accept()
    print(f"[{dlp.name}] Sender bağlandı:", sender_addr)

    sender_file = sender_sock.makefile("r", encoding="utf-8")

    try:
        for line in sender_file:
            text = line.rstrip("\n")
            if not text:
                continue

            msg = Message(
                src="SENDER_PC",
                dst="RECEIVER_PC",
                channel="chat",
                payload=text,
            )

            checked = dlp.handle(msg)

            if checked is None:
                sender_sock.sendall(
                    "[DLP] Mesajın hassas veri içerdiği için gönderilmedi.\n".encode("utf-8")
                )
            else:
                receiver_sock.sendall((checked.payload + "\n").encode("utf-8"))

    except KeyboardInterrupt:
        print("\n[Gateway] Kapatılıyor...")
    finally:
        sender_file.close()
        sender_sock.close()
        receiver_sock.close()
        server_sock.close()


def run_sender():
    print("========================================")
    print("   Endpoint Sender Agent (PC1)         ")
    print("========================================")
    print(f"DLP Gateway: {SENDER_GATEWAY_HOST}:{SENDER_GATEWAY_PORT}")
    print("Çıkmak için 'q' yaz.\n")

    sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    sock.connect((SENDER_GATEWAY_HOST, SENDER_GATEWAY_PORT))
    print("[SENDER] Gateway'e bağlanıldı.\n")

    gateway_file = sock.makefile("r", encoding="utf-8")

    try:
        while True:
            text = input("Gönderilecek mesaj: ").strip()
            if text.lower() in {"q", "quit", "exit"}:
                print("[SENDER] Çıkılıyor...")
                break

            if not text:
                continue

            sock.sendall((text + "\n").encode("utf-8"))

            sock.settimeout(0.2)
            try:
                line = gateway_file.readline()
                if line:
                    print("[GATEWAY MESAJI]", line.strip())
            except Exception:
                pass
            finally:
                sock.settimeout(None)

    except KeyboardInterrupt:
        print("\n[SENDER] Kapatılıyor...")
    finally:
        gateway_file.close()
        sock.close()


def run_receiver():
    print("========================================")
    print("   Endpoint Receiver Agent (PC2)       ")
    print("========================================")
    print(f"{RECEIVER_LISTEN_HOST}:{RECEIVER_LISTEN_PORT} dinleniyor...\n")

    server_sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    server_sock.bind((RECEIVER_LISTEN_HOST, RECEIVER_LISTEN_PORT))
    server_sock.listen(1)

    conn, addr = server_sock.accept()
    print("[RECEIVER] Gateway bağlandı:", addr)

    conn_file = conn.makefile("r", encoding="utf-8")

    try:
        for line in conn_file:
            text = line.rstrip("\n")
            if not text:
                continue
            print(f"[RECEIVER] Yeni mesaj: {text}")
    except KeyboardInterrupt:
        print("\n[RECEIVER] Kapatılıyor...")
    finally:
        conn_file.close()
        conn.close()
        server_sock.close()

def fetch_policy():
    global USB_POLICY, USB_SCAN_ON_MODIFY
    global CLIPBOARD_ENABLED, USB_ENABLED, NETWORK_ENABLED
    global BANNED_WORDS, SCAN_RULES

    try:
        r = requests.get(POLICY_URL, timeout=3).json()

        if "error" in r:
            print("[Policy] Sunucu 'unknown endpoint' dedi!")
            return False

        # --- USB & Features ---
        USB_POLICY = r.get("usb_policy", USB_POLICY)
        USB_SCAN_ON_MODIFY = r.get("scan_on_modify", True)
        features = r.get("features", {})
        CLIPBOARD_ENABLED = features.get("clipboard_enabled", CLIPBOARD_ENABLED)
        USB_ENABLED       = features.get("usb_enabled", USB_ENABLED)
        NETWORK_ENABLED   = features.get("network_enabled", NETWORK_ENABLED)

        # --- banned words / scan rules ---
        BANNED_WORDS = set(map(str.lower, r.get("banned_words", [])))
        if "scan_rules" in r:
            SCAN_RULES.update(r["scan_rules"])

        # Kalıcı olarak kaydet
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump({
                "USB_POLICY": USB_POLICY,
                "SCAN_RULES": SCAN_RULES,
                "BANNED_WORDS": list(BANNED_WORDS)
            }, f, indent=4)
        print("[Config] dlp_config.json kalıcı olarak güncellendi.")

        print("\n[Policy] Policy başarıyla güncellendi!")
        print(f"USB_POLICY        = {USB_POLICY}")
        print(f"CLP:{CLIPBOARD_ENABLED} | USB:{USB_ENABLED} | NET:{NETWORK_ENABLED}")
        print(f"SCAN_RULES        = {SCAN_RULES}")
        print(f"BANNED_WORDS      = {BANNED_WORDS}")

        return True

    except Exception as e:
        print("[Policy] Policy çekilemedi:", e)
        return False

    
def update_policy_on_server(new_policy):
    """
    Endpoint'teki değişiklikleri Policy Server'a gönderir (POST).
    new_policy: dict formatında policy bilgisi
    """
    try:
        url = "http://127.0.0.1:5000/update_policy"  # server route
        r = requests.post(url, json=new_policy, timeout=3)

        # Sunucunun cevabını göster
        print("\n[Policy] Sunucudan gelen cevap:")
        try:
            print(json.dumps(r.json(), indent=4, ensure_ascii=False))
        except:
            print(r.text)

    except Exception as e:
        print("[Policy ERROR]:", e)

def print_current_policy():
    print("\n=== AKTİF POLICY (Endpoint Üzerinde) ===")
    print(f"USB_POLICY        : {USB_POLICY}")
    print(f"USB_SCAN_ON_MODIFY: {USB_SCAN_ON_MODIFY}")
    print(f"CLIPBOARD_ENABLED : {CLIPBOARD_ENABLED}")
    print(f"USB_ENABLED       : {USB_ENABLED}")
    print(f"NETWORK_ENABLED   : {NETWORK_ENABLED}")
    print(f"BANNED_WORDS      : {list(BANNED_WORDS)}")
    print(f"SCAN_RULES        : {SCAN_RULES}")
    print("========================================\n")



# ============================================================
# MAIN / MENÜ
# ============================================================

def run_endpoint_dlp():
    usb_thread = threading.Thread(target=usb_monitor, daemon=True)
    usb_thread.start()
    clipboard_thread = threading.Thread(target=clipboard_monitor, daemon=True)
    clipboard_thread.start()

    print("---------------------------------------------------------")
    print("--- Mini DLP Endpoint Modu Başlatıldı ---")
    print("Pano Kontrolü: Aktif")
    print("USB Kontrolü: Aktif (Gerçek mount noktaları taranıyor)")
    print("Kapsam: TCKN, Telefon, Kredi Kartı, E-posta, IBAN, Yasaklı Kelimeler")
    print("Durdurmak için: CTRL+C\n")

    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        print("\n--- DLP Endpoint Modu Durduruluyor ---")
        print("Durduruldu. Olaylar 'dlp_incidents.csv' dosyasında kayıtlıdır.")


def main_menu():
    
    print(f"Aktif USB Policy: {USB_POLICY}  (server veya config.json üzerinden değiştirilebilir)")

    print("===================================================")
    print("   Tek Dosyalık DLP Agent Sistemi (Mini Proje)     ")
    print("===================================================")
    print("Mod Seç:")
    print("  0) Policy Server'dan policy çek (GET)")
    print("  1) Endpoint DLP (Clipboard + USB izle)")
    print("  2) Sender Agent (PC1)")
    print("  3) DLP Gateway (Aradaki AI/DLP Agent)")
    print("  4) Receiver Agent (PC2)")
    print("  5) Yasaklı kelime ekle")
    print("  6) Yasaklı kelime kaldır")
    print("  7) Yasaklı kelimeleri listele")
    print("  8) Tarama ayarlarını değiştir (SCAN_RULES)")
    print("  9) Aktif policy'yi göster")
    print(" 10) Policy sunucuya gönder (POST)")
    print("  q) Çık")
    print("===================================================\n")

    choice = input("Seçimin: ").strip().lower()

    if choice == "0":
        fetch_policy()
        return main_menu()
    if choice == "1":
        if CLIPBOARD_ENABLED:
            start_clipboard_monitor()
        if USB_ENABLED:
            start_usb_monitor()
        if NETWORK_ENABLED:
            start_network_monitor()  # varsa
        print("\n--- DLP Endpoint Başlatıldı ---")
        print(f"Clipboard: {CLIPBOARD_ENABLED}")
        print(f"USB:       {USB_ENABLED} ({USB_POLICY})")
        print(f"Network:   {NETWORK_ENABLED}\n")
    elif choice == "2":
        run_sender()
    elif choice == "3":
        run_gateway()
    elif choice == "4":
        run_receiver()
    elif choice == "5":
        word = input("Eklenecek kelime: ").strip()
        add_banned_word(word)
        main_menu()
    elif choice == "6":
        word = input("Kaldırılacak kelime: ").strip()
        remove_banned_word(word)
        main_menu()
    elif choice == "7":
        list_banned_words()
        main_menu()
    elif choice == "8":
        print("\n--- Tarama Ayarları (SCAN_RULES) ---")
        print("Mevcut Ayarlar:")
        for rule, status in SCAN_RULES.items():
            print(f" - {rule}: {'AKTIF' if status else 'PASIF'}")

        rule = input("\nDeğiştirilecek veri tipi (örn: IBAN_TR) / q: geri çık: ").strip()

        if rule in SCAN_RULES:
            SCAN_RULES[rule] = not SCAN_RULES[rule]
            print(f"{rule} artık {'AKTIF' if SCAN_RULES[rule] else 'PASIF'}")

            # değişiklik sonrası json’a yaz
            with open(CONFIG_FILE, "w", encoding="utf-8") as f:
                json.dump({"SCAN_RULES": SCAN_RULES}, f, indent=4)
            print("[Config] dlp_config.json güncellendi.")

        main_menu()
    elif choice == "9":
        print_current_policy()
        main_menu()
    elif choice == "10":
        new_policy = {
            "usb_policy": USB_POLICY,
            "scan_on_modify": USB_SCAN_ON_MODIFY,
            "features": {
                "clipboard_enabled": CLIPBOARD_ENABLED,
                "usb_enabled": USB_ENABLED,
                "network_enabled": NETWORK_ENABLED
            },
            "banned_words": list(BANNED_WORDS),
            "scan_rules": SCAN_RULES
        }
        update_policy_on_server(new_policy)
        main_menu()
    elif choice in {"q", "quit", "exit"}:
        print("Çıkılıyor...")
    else:
        print("Geçersiz seçim.")


if __name__ == "__main__":
    main_menu()
