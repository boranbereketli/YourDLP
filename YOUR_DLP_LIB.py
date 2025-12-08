# YOUR_DLP_LIB.py

import re
import os
import shutil
import time
from dataclasses import dataclass
from typing import Optional
import pandas as pd
from docx import Document
import PyPDF2
from pptx import Presentation
import string
import ctypes
import platform

# ============================================================
# CONFIG & AYARLAR
# ============================================================

LOG_CSV = "dlp_incidents.csv"
QUARANTINE_DIR = "KARANTINA_ALANI"
MAX_FILE_SIZE = 15 * 1024 * 1024  # 15 MB
ALLOWED_EXT = {".txt", ".csv", ".docx", ".pdf", ".xlsx", ".xls", ".pptx"}
DLP_SCAN_ORDER = ["TCKN", "TEL_NO", "IBAN_TR", "KREDI_KARTI", "E_POSTA"]

# ------------------------------------------------------------
# Regex Patterns (Pre-compiled for performance)
# ------------------------------------------------------------
# 🚨 ÖNEMLİ: \b (kelime sınırı) karakterleri, parça tespiti için kaldırılmıştır!

# TCKN: Sadece 11 haneli bitişik rakamları arar.
REGEX_TCKN = re.compile(r'\d{11}') 

# TEL_NO: Aynı kalır, zaten \b ihtiyacı yoktur.
REGEX_TEL = re.compile(r'(?:(?:\+90|0)?5\d{9})')

# KREDİ KARTI: 16 haneli formatı arar. (Önceki: r'\b... \b')
REGEX_CC = re.compile(r'\d{4}[- ]?\d{4}[- ]?\d{4}[- ]?\d{4}')

# E-POSTA: Aynı kalır.
REGEX_EMAIL = re.compile(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}')

# IBAN: TR ile başlayan formatı arar. (Önceki: r'\b... \b')
REGEX_IBAN = re.compile(r'TR\d{2}[A-Z0-9]{4}\s?(?:\d{4}\s?){4}\d{2}')


# ------------------------------------------------------------
# DLP_RULES Sözlüğü
# ------------------------------------------------------------
DLP_RULES = {
    "TCKN": {"pattern": REGEX_TCKN, "description": "11 Haneli TC Kimlik Numarası"},
    "TEL_NO": {"pattern": REGEX_TEL, "description": "Türkiye Telefon Numarası"},
    "KREDI_KARTI": {"pattern": REGEX_CC, "description": "16 Haneli Kredi Kartı Numarası Formatı"},
    "E_POSTA": {"pattern": REGEX_EMAIL, "description": "E-posta Adresi Formatı"},
    "IBAN_TR": {"pattern": REGEX_IBAN, "description": "Türk IBAN Numarası Formatı"}
}


@dataclass
class Message:
    src: str
    dst: str
    channel: str
    payload: str


# ============================================================
# HELPERS: VALIDATORS
# ============================================================

def is_valid_tckn(tckn: str) -> bool:
    if not isinstance(tckn, str): return False
    t = re.sub(r"\D", "", tckn)
    if len(t) != 11 or t[0] == "0": return False
    try:
        digits = list(map(int, t))
    except ValueError:
        return False

    rule_10 = ((digits[0] + digits[2] + digits[4] + digits[6] + digits[8]) * 7 -
               (digits[1] + digits[3] + digits[5] + digits[7])) % 10
    if rule_10 != digits[9]: return False

    rule_11 = sum(digits[:10]) % 10
    if rule_11 != digits[10]: return False
    return True


def is_valid_phone(phone: str) -> bool:
    if not isinstance(phone, str): return False
    digits = re.sub(r"\D", "", phone)
    
    if digits.startswith("90") and len(digits) == 12:
        national = digits[2:]
    elif digits.startswith("0") and len(digits) == 11:
        national = digits[1:]
    elif len(digits) == 10:
        national = digits
    else:
        return False

    if len(national) != 10 or not national.isdigit() or not national.startswith("5"):
        return False
    return True


def iban_to_numeric(iban: str) -> Optional[str]:
    try:
        s = re.sub(r"\s+", "", iban).upper()
        if len(s) < 4: return None
        rearr = s[4:] + s[:4]
        numeric = []
        for ch in rearr:
            if ch.isdigit():
                numeric.append(ch)
            elif ch.isalpha():
                numeric.append(str(ord(ch) - 55))
            else:
                return None
        return "".join(numeric)
    except Exception:
        return None


def is_valid_iban(iban: str) -> bool:
    if not isinstance(iban, str): return False
    s = re.sub(r"\s+", "", iban).upper()
    if not re.match(r"^[A-Z]{2}\d{2}[A-Z0-9]+$", s): return False
    numeric = iban_to_numeric(s)
    if not numeric: return False

    remainder = 0
    chunk_size = 9
    for i in range(0, len(numeric), chunk_size):
        chunk = str(remainder) + numeric[i:i + chunk_size]
        remainder = int(chunk) % 97
    return remainder == 1


# ============================================================
# SCANNER CORE
# ============================================================

def scan_content(content: str, dynamic_keywords: list = None):
    """ 
    Tüm hassas veri tiplerini ve dinamik anahtar kelimeleri tarar. 
    Kelime sınırlarını göz ardı ederek parça tespiti (substring matching) yapar.
    """
    incidents = []
    if not content: return incidents
    try: full_text = str(content)
    except Exception: full_text = ""

    upper_text = full_text.upper() 

    # --- 1) DİNAMİK ANAHTAR KELİME TARAMASI ---
    if dynamic_keywords:
        for keyword in dynamic_keywords:
            upper_keyword = keyword.upper()
            if upper_keyword in upper_text:
                # Anahtar kelimeyi 'KEYWORD_MATCH' olarak kaydet
                incidents.append({
                    "data_type": "KEYWORD_MATCH",
                    "description": f"Anahtar Kelime Tespiti: {keyword}",
                    "masked_match": f"[KEYWORD] {keyword[:15]}..."
                })
                # Bir eşleşme bulmak yeterli (listeden çıkan kurala göre)
                # NOT: Tüm listeyi kontrol etmek yerine, bu veri tipini engellemek yeterlidir.
                break 

    # --- 2) REGEX ve NUMERİK TARAMA (Parça Tespiti Dahil) ---
    text_for_tel_no = list(full_text)
    
    # A) TCKN Tespiti (Agresif 11 haneli sayı blokları arama)
    
    # 11 haneli veya daha uzun bitişik rakam dizilerini arar
    tckn_candidate_pattern = re.compile(r'\d{11,}') 
    tckn_matches = set()
    
    # TCKN için kayan pencere kontrolü
    for mo in tckn_candidate_pattern.finditer(full_text):
        full_match = mo.group(0)
        # 11 hanelik tüm alt dizileri kontrol et (kayan pencere)
        for i in range(len(full_match) - 10):
            cand = full_match[i:i+11]
            if is_valid_tckn(cand):
                tckn_matches.add(cand)
                
                # Doğrulanan TCKN'leri, telefon taramasından çıkarmak için metinde boşlukla değiştir.
                start_index = mo.span()[0] + i 
                for j in range(11):
                     idx = start_index + j
                     if 0 <= idx < len(text_for_tel_no):
                         text_for_tel_no[idx] = " "
    
    # TCKN Incident'larını ekle
    rule_tckn = DLP_RULES["TCKN"]
    for cand in sorted(list(tckn_matches)):
        masked = f"TC: ******{cand[-4:]}"
        incidents.append({"data_type": "TCKN", "description": rule_tckn["description"], "masked_match": masked})

    text_for_tel_no = "".join(text_for_tel_no) # Telefon taramasına hazır metin

    # B) Telefon Numaraları (TCKN çıkarılmış metin üzerinde arama)
    rule_tel = DLP_RULES["TEL_NO"]
    try: tel_matches = re.findall(rule_tel["pattern"], text_for_tel_no)
    except Exception: tel_matches = []

    for m in set(tel_matches):
        if is_valid_phone(m):
            flat = re.sub(r"\D", "", m)
            masked = f"TEL: ******{flat[-2:]}"
            incidents.append({"data_type": "TEL_NO", "description": rule_tel["description"], "masked_match": masked})

    # C) Diğer kurallar (IBAN_TR, KREDI_KARTI, E_POSTA)
    for data_type in DLP_SCAN_ORDER:
        if data_type in {"TCKN", "TEL_NO"}: continue
        
        rule = DLP_RULES[data_type]
        # Regex'ler artık \b olmadan tanımlandığı için kelime parçalarını yakalar
        try: matches = re.findall(rule["pattern"], full_text) 
        except re.error: matches = []

        for match in set(matches):
            if isinstance(match, tuple): match = "".join(match)

            if data_type == "IBAN_TR":
                cand = re.sub(r"\s+", "", match).upper()
                if is_valid_iban(cand):
                    masked = f"IBAN: ****{cand[-4:]}"
                    incidents.append({"data_type": "IBAN_TR", "description": rule["description"], "masked_match": masked})

            elif data_type == "KREDI_KARTI":
                # Kredi kartı formatını (16 rakam) kelimeye bitişik bulsa bile kaydeder
                flat = re.sub(r"\D", "", match)
                if len(flat) == 16: # Format eksiksiz olmalı
                    masked = f"CC: XXXX...{flat[-4:]}"
                    incidents.append({"data_type": "KREDI_KARTI", "description": rule["description"], "masked_match": masked})

            elif data_type == "E_POSTA":
                try:
                    name_part = match.split('@')[0]
                    masked = f"EMAIL: <{name_part[0]}***@...>"
                    incidents.append({"data_type": "E_POSTA", "description": rule["description"], "masked_match": masked})
                except: continue

    return incidents


# ============================================================
# FILE READING
# ============================================================

def read_file_content(path: str) -> str:
    """ Dosya içeriğini güvenli şekilde okur. """
    if not os.path.exists(path): return ""
    try:
        if os.path.getsize(path) > MAX_FILE_SIZE: return ""
        
        ext = os.path.splitext(path)[1].lower()

        if ext in (".txt", ".csv"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
                
        if ext == ".docx":
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
            
        if ext == ".pdf":
            text = ""
            with open(path, "rb") as f:
                try:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() or ""
                except: pass
            return text
            
        if ext in (".xlsx", ".xls"):
            try:
                dfs = pd.read_excel(path, sheet_name=None)
                return "\n".join(df.to_string(index=False) for df in dfs.values())
            except: return ""
            
        if ext == ".pptx":
            try:
                prs = Presentation(path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            parts.append(shape.text)
                return "\n".join(parts)
            except: return ""

        return ""
    except Exception as e:
        print(f"[READ ERROR] {e}")
        return ""


# ============================================================
# USB & SYSTEM HELPERS
# ============================================================

def get_usb_mount_points(sim_usb_dir):
    mounts = []
    sys_name = platform.system()
    
    try:
        if sys_name == "Windows":
            DRIVE_REMOVABLE = 2
            bitmask = ctypes.windll.kernel32.GetLogicalDrives()
            for i, letter in enumerate(string.ascii_uppercase):
                if bitmask & (1 << i):
                    drive = f"{letter}:\\"
                    try:
                        if ctypes.windll.kernel32.GetDriveTypeW(ctypes.c_wchar_p(drive)) == DRIVE_REMOVABLE:
                            mounts.append(drive)
                    except: continue
        else:
            # Linux/Mac
            base_dirs = ["/media", "/run/media", "/mnt"]
            user = os.getenv("USER") or "root"
            if user != "root":
                base_dirs.extend([f"/run/media/{user}", f"/media/{user}"])
            
            for base in base_dirs:
                if os.path.exists(base):
                    for entry in os.listdir(base):
                        candidate = os.path.join(base, entry)
                        if os.path.isdir(candidate) and not os.path.islink(candidate):
                            mounts.append(candidate)
    except: pass

    if os.path.exists(sim_usb_dir) and sim_usb_dir not in mounts:
        mounts.append(sim_usb_dir)
    return list(set(mounts))

def quarantine_file(src_path, quarantine_dir=QUARANTINE_DIR, hint_name=None):
    name = hint_name or os.path.basename(src_path)
    dest = os.path.join(quarantine_dir, f"{int(time.time())}_{name}")
    os.makedirs(quarantine_dir, exist_ok=True)
    try:
        shutil.move(src_path, dest)
        return dest
    except Exception:
        # Move başarısızsa copy+delete dene
        try:
            shutil.copy2(src_path, dest)
            os.remove(src_path)
            return dest
        except: return None
