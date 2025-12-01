# YOUR_DLP_LIB.py

import re
import os
import shutil
import time
from dataclasses import dataclass
from typing import Optional
import pandas as pd
from docx import Document
import PyPDF2
from pptx import Presentation
import string
import ctypes
import platform

# ============================================================
# CONFIG & AYARLAR
# ============================================================

LOG_CSV = "dlp_incidents.csv"
QUARANTINE_DIR = "KARANTINA_ALANI"
MAX_FILE_SIZE = 15 * 1024 * 1024  # 15 MB
ALLOWED_EXT = {".txt", ".csv", ".docx", ".pdf", ".xlsx", ".xls", ".pptx"}
DLP_SCAN_ORDER = ["TCKN", "TEL_NO", "IBAN_TR", "KREDI_KARTI", "E_POSTA"]

# ============================================================
# DATA STRUCTURES
# ============================================================

@dataclass
class Message:
    src: str
    dst: str
    channel: str
    payload: str


# ============================================================
# HELPERS: TCKN / PHONE / IBAN VALIDATORS (Aynı Kalır)
# ============================================================

def is_valid_tckn(tckn: str) -> bool:
    """ Türkiye Cumhuriyeti Kimlik Numarası doğrulaması """
    if not isinstance(tckn, str): return False
    t = re.sub(r"\D", "", tckn)
    if len(t) != 11 or t[0] == "0": return False
    try: digits = list(map(int, t))
    except ValueError: return False

    rule_10 = ((digits[0] + digits[2] + digits[4] + digits[6] + digits[8]) * 7 - 
               (digits[1] + digits[3] + digits[5] + digits[7])) % 10
    if rule_10 != digits[9]: return False

    rule_11 = sum(digits[:10]) % 10
    if rule_11 != digits[10]: return False
    return True

def is_valid_phone(phone: str) -> bool:
    """ Türkiye mobil telefon doğrulaması """
    if not isinstance(phone, str): return False
    digits = re.sub(r"\D", "", phone)
    if digits.startswith("90") and len(digits) == 12: national = digits[2:]
    elif digits.startswith("0") and len(digits) == 11: national = digits[1:]
    elif len(digits) == 10: national = digits
    else: return False

    if len(national) != 10 or not national.isdigit() or not national.startswith("5"):
        return False
    return True

def iban_to_numeric(iban: str) -> Optional[str]:
    """ IBAN string -> numeric string for mod-97 """
    try:
        s = re.sub(r"\s+", "", iban).upper()
        if len(s) < 4: return None
        rearr = s[4:] + s[:4]
        numeric = []
        for ch in rearr:
            if ch.isdigit(): numeric.append(ch)
            elif ch.isalpha(): numeric.append(str(ord(ch) - 55))
            else: return None
        return "".join(numeric)
    except Exception: return None

def is_valid_iban(iban: str) -> bool:
    """ IBAN mod-97 validation. """
    if not isinstance(iban, str): return False
    s = re.sub(r"\s+", "", iban).upper()
    if not re.match(r"^[A-Z]{2}\d{2}[A-Z0-9]+$", s): return False
    numeric = iban_to_numeric(s)
    if not numeric: return False
    
    remainder = 0
    chunk_size = 9
    for i in range(0, len(numeric), chunk_size):
        chunk = str(remainder) + numeric[i:i + chunk_size]
        remainder = int(chunk) % 97
    return remainder == 1

# ============================================================
# DLP RULES & SCANNER 
# ============================================================

TEL_REGEX = r'(?:(?:\+90|0)?5\d{9})' 
DLP_RULES = {
    "TCKN": {"pattern": r'\b[1-9]\d{10}\b', "description": "11 Haneli TC Kimlik Numarası"}, 
    "TEL_NO": {"pattern": TEL_REGEX, "description": "Türkiye Telefon Numarası"},
    "KREDI_KARTI": {"pattern": r'\b\d{4}[- ]?\d{4}[- ]?\d{4}[- ]?\d{4}\b', "description": "16 Haneli Kredi Kartı Numarası Formatı"},
    "E_POSTA": {"pattern": r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', "description": "E-posta Adresi Formatı"},
    "IBAN_TR": {"pattern": r'\bTR\d{2}[A-Z0-9]{4}\s?(?:\d{4}\s?){4}\d{2}\b', "description": "Türk IBAN Numarası Formatı"}
}


def scan_content(content: str, dynamic_keywords: list = None):
    """ Tüm hassas veri tiplerini tarar ve bulunanları listeler. """
    incidents = []
    if not content: return incidents
    try: full_text = str(content)
    except Exception: full_text = ""
    
    upper_text = full_text.upper() 

    # --- 1) DİNAMİK ANAHTAR KELİME TARAMASI ---
    if dynamic_keywords:
        for keyword in dynamic_keywords:
            upper_keyword = keyword.upper()
            if upper_keyword in upper_text:
                # Anahtar kelimeyi 'KEYWORD_MATCH' olarak kaydet
                incidents.append({
                    "data_type": "KEYWORD_MATCH",
                    "description": f"Anahtar Kelime Tespiti: {keyword}",
                    "masked_match": f"[KEYWORD] {keyword[:15]}..."
                })


    text_for_tel_no = list(full_text)
    
    # 1) TCKN'leri tespit et, doğrula ve metinden çıkar
    rule = DLP_RULES["TCKN"]
    try: tckn_pattern = re.compile(rule["pattern"])
    except re.error: matches = []
    else: matches = tckn_pattern.finditer(full_text)
    
    tckn_matches = set()
    for mo in matches:
        match = mo.group(0)
        cand = re.sub(r"\D", "", match)
        if is_valid_tckn(cand):
            tckn_matches.add(cand)
            start, end = mo.span()
            for i in range(start, end):
                if 0 <= i < len(text_for_tel_no):
                    text_for_tel_no[i] = " "
    
    for cand in sorted(list(tckn_matches)):
        masked = f"TC: ******{cand[-4:]}"
        incidents.append({"data_type": "TCKN", "description": rule["description"], "masked_match": masked})

    text_for_tel_no = "".join(text_for_tel_no)

    # 2) Telefon Numaralarını tespit et
    rule = DLP_RULES["TEL_NO"]
    try: tel_matches = re.findall(rule["pattern"], text_for_tel_no)
    except Exception: tel_matches = []

    for m in set(tel_matches):
        if is_valid_phone(m):
            flat = re.sub(r"\D", "", m)
            masked = f"TEL: ******{flat[-2:]}"
            incidents.append({"data_type": "TEL_NO", "description": rule["description"], "masked_match": masked})

    # 3) Diğer kurallar (IBAN_TR, KREDI_KARTI, E_POSTA)
    for data_type in DLP_SCAN_ORDER:
        if data_type in {"TCKN", "TEL_NO"}: continue
        rule = DLP_RULES[data_type]
        try: matches = re.findall(rule["pattern"], full_text)
        except re.error: matches = []

        for match in set(matches):
            if isinstance(match, tuple): match = "".join(match)

            if data_type == "IBAN_TR":
                cand = re.sub(r"\s+", "", match).upper()
                if is_valid_iban(cand):
                    masked = f"IBAN: ****{cand[-4:]}"
                    incidents.append({"data_type": "IBAN_TR", "description": rule["description"], "masked_match": masked})

            elif data_type == "KREDI_KARTI":
                flat = re.sub(r"\D", "", match)
                masked = f"CC: XXXX...{flat[-4:]}"
                incidents.append({"data_type": "KREDI_KARTI", "description": rule["description"], "masked_match": masked})

            elif data_type == "E_POSTA":
                try:
                    name_part = match.split('@')[0]
                    masked = f"EMAIL: <{name_part[0]}***@...>"
                    incidents.append({"data_type": "E_POSTA", "description": rule["description"], "masked_match": masked})
                except: continue

    return incidents


# ============================================================
# FILE READING (Aynı Kalır)
# ============================================================

def read_file_content(path: str) -> str:
    """ Dosya içeriğini okur (TXT, DOCX, PDF, XLSX, PPTX) """
    try: ext = os.path.splitext(path)[1].lower()
    except Exception: return ""

    try:
        if os.path.getsize(path) > MAX_FILE_SIZE: return ""

        if ext in (".txt", ".csv"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f: return f.read()
        if ext == ".docx":
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        if ext == ".pdf":
            text = ""
            with open(path, "rb") as f:
                try:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages: text += page.extract_text() or ""
                except Exception: return ""
            return text
        if ext in (".xlsx", ".xls"):
            try:
                dfs = pd.read_excel(path, sheet_name=None)
                return "\n".join(df.to_string(index=False) for df in dfs.values())
            except Exception: return ""
        if ext == ".pptx":
            try:
                prs = Presentation(path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): parts.append(shape.text)
                return "\n".join(parts)
            except Exception: return ""
        return ""
    except Exception as e:
        print(f"[READ ERROR] {e}")
        return ""


# ============================================================
# USB MOUNT POINT & QUARANTINE HELPERS (Aynı Kalır)
# ============================================================

def get_usb_mount_points(sim_usb_dir):
    """ İşletim sistemine göre USB mount noktalarını bulur. """
    mounts = []
    system = platform.system()
    try:
        if system == "Windows":
            DRIVE_REMOVABLE = 2
            bitmask = ctypes.windll.kernel32.GetLogicalDrives()
            for i, letter in enumerate(string.ascii_uppercase):
                if bitmask & (1 << i):
                    drive = f"{letter}:\\"
                    try:
                        drive_type = ctypes.windll.kernel32.GetDriveTypeW(ctypes.c_wchar_p(drive))
                        if drive_type == DRIVE_REMOVABLE: mounts.append(drive)
                    except Exception: continue
        else:
            base_dirs = ["/media", "/run/media", "/mnt"]
            user = os.getenv("USER") or "root"
            if user != "root": base_dirs.extend([f"/run/media/{user}", f"/media/{user}"])
            for base in base_dirs:
                if os.path.exists(base):
                    for entry in os.listdir(base):
                        candidate = os.path.join(base, entry)
                        if os.path.isdir(candidate) and not os.path.islink(candidate): mounts.append(candidate)
                            
    except Exception: pass

    if os.path.exists(sim_usb_dir) and sim_usb_dir not in mounts:
        mounts.append(sim_usb_dir)
    return list(set(mounts))

def quarantine_file(src_path, quarantine_dir=QUARANTINE_DIR, hint_name=None):
    """ Dosyayı karantinaya taşır. """
    name = hint_name or os.path.basename(src_path)
    dest = os.path.join(quarantine_dir, f"{int(time.time())}_{name}")
    os.makedirs(quarantine_dir, exist_ok=True)
    try:
        shutil.move(src_path, dest)
        return dest
    except Exception:
        try:
            shutil.copy2(src_path, dest)
            try: os.remove(src_path)
            except Exception: pass
            return dest
        except Exception: return None
